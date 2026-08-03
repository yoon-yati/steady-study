from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


def read_txt(uploaded_file):

    return uploaded_file.read().decode(
        "utf-8",
        errors="ignore"
    )


def read_pdf(uploaded_file):

    text = ""

    try:

        pdf = PdfReader(uploaded_file)

        for page in pdf.pages:

            page_text = page.extract_text()

            if page_text:

                text += page_text + "\n"

    except Exception:

        return ""

    return text


def read_docx(uploaded_file):

    try:

        doc = Document(uploaded_file)

        text = "\n".join(
            paragraph.text
            for paragraph in doc.paragraphs
        )

        return text

    except Exception:

        return ""


def read_pptx(uploaded_file):

    text = ""

    try:

        prs = Presentation(uploaded_file)

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

    except Exception:

        return ""

    return text


def extract_text(uploaded_file):

    filename = uploaded_file.name.lower()

    if filename.endswith(".txt"):

        return read_txt(
            uploaded_file
        )

    elif filename.endswith(".pdf"):

        return read_pdf(
            uploaded_file
        )

    elif filename.endswith(".docx"):

        return read_docx(
            uploaded_file
        )

    elif filename.endswith(".pptx"):

        return read_pptx(
            uploaded_file
        )

    return ""